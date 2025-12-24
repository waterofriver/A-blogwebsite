from django.shortcuts import render, get_object_or_404
from django.http import FileResponse, Http404, HttpResponse
from django.utils import timezone
from django.db.models import Q
from django.conf import settings
from django.urls import reverse
from rest_framework import viewsets, permissions, status, filters
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.views import APIView
from django_filters.rest_framework import DjangoFilterBackend
from .models import MaterialCategory, MaterialItem, Attachment
from .serializers import (
    MaterialCategorySerializer, MaterialCategoryListSerializer,
    MaterialItemSerializer, MaterialItemListSerializer, 
    AttachmentSerializer
)
from pathlib import Path, PurePosixPath
from urllib.parse import quote
from datetime import datetime, timezone as dt_timezone
from typing import Optional
import os
import mimetypes
import json
import html

from docx import Document
from pptx import Presentation

MATERIALS_FILENAME = 'materials.json'
MATERIALS_SUBDIR = 'materials'
BOOK_CATEGORIES = {'参考书籍', '参考资料'}
VIDEO_FILE_TYPES = {'mp4', 'mov', 'avi', 'wmv', 'mkv'}
BOOK_CATEGORY_ORDER = ['参考书籍', '参考资料']


def _materials_json_path() -> Path:
    return Path(settings.BASE_DIR) / 'data' / MATERIALS_FILENAME


def _extract_category_order(category: Optional[str]) -> int:
    if not category:
        return 999
    prefix = category.split('_', 1)[0]
    return int(prefix) if prefix.isdigit() else 999


def _category_display(category: Optional[str]) -> str:
    if not category:
        return '未分类'
    parts = category.split('_', 1)
    if len(parts) == 2 and parts[0].isdigit():
        return parts[1].strip() or category
    return category


def _normalize_material_path(relative_path: Optional[str]) -> str:
    if not relative_path:
        raise ValueError('Empty path')
    clean = relative_path.strip().lstrip('/\\')
    if not clean:
        raise ValueError('Empty path')
    pure = PurePosixPath(clean)
    if pure.is_absolute() or '..' in pure.parts:
        raise ValueError('Invalid path')
    return pure.as_posix()


def _resolve_material_file(relative_path: str) -> tuple[str, Path]:
    """Return the storage-relative path (posix) and absolute Path for the given entry."""
    normalized = _normalize_material_path(relative_path)
    base_dir = Path(settings.MEDIA_ROOT) / MATERIALS_SUBDIR
    pure = PurePosixPath(normalized)

    def _path_for(candidate: PurePosixPath) -> Path:
        return base_dir / Path(*candidate.parts)

    direct_path = _path_for(pure)
    if direct_path.exists() and direct_path.is_file():
        return pure.as_posix(), direct_path

    # try simple fallback to filename at root
    filename_only = PurePosixPath(pure.name)
    root_path = _path_for(filename_only)
    if root_path.exists() and root_path.is_file():
        return filename_only.as_posix(), root_path

    # search entire tree for the filename
    try:
        matched = next(base_dir.rglob(pure.name))
    except StopIteration:
        matched = None

    if matched and matched.is_file():
        relative = matched.relative_to(base_dir).as_posix()
        return relative, matched

    raise FileNotFoundError(f"{relative_path} not found in {base_dir}")


INLINE_PREVIEW_TYPES = {
    'pdf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'mp4', 'mov', 'avi', 'wmv', 'mkv'
}

DOC_EXTENSIONS = {'.doc', '.docx'}
PPT_EXTENSIONS = {'.ppt', '.pptx'}
TEXT_EXTENSIONS = {'.txt', '.md'}


def _build_attachment_payload(entry: dict, request) -> dict:
    storage_relative, _ = _resolve_material_file(entry.get('filename'))
    category = entry.get('category') or '未分类'
    category_order = _extract_category_order(category)
    category_display = _category_display(category)
    label = entry.get('label') or entry.get('orig_name') or PurePosixPath(storage_relative).name
    encoded_path = quote(storage_relative, safe='/')
    media_url = request.build_absolute_uri(f"{settings.MEDIA_URL}{MATERIALS_SUBDIR}/{encoded_path}")
    preview_url = request.build_absolute_uri(
        f"{reverse('material-file-preview')}?path={quote(storage_relative)}"
    )
    download_url = request.build_absolute_uri(
        f"{reverse('material-file-download')}?path={quote(storage_relative)}"
    )
    file_type = (entry.get('file_type') or '').lower()
    return {
        'id': f"{category}:{storage_relative}",
        'category': category,
        'category_display': category_display,
        'category_order': category_order,
        'item_name': entry.get('item_name'),
        'item_label': entry.get('item_label'),
        'label': label,
        'filename': storage_relative,
        'file_type': file_type,
        'media_url': media_url,
        'preview_url': preview_url,
        'download_url': download_url,
        'html_preview_url': (
            request.build_absolute_uri(
                f"{reverse('material-file-html-preview')}?path={quote(storage_relative)}"
            )
            if PurePosixPath(storage_relative).suffix.lower() in DOC_EXTENSIONS | PPT_EXTENSIONS | TEXT_EXTENSIONS
            else None
        ),
        'orig_name': entry.get('orig_name'),
        'supports_inline_preview': file_type in INLINE_PREVIEW_TYPES,
    }


def _render_docx_to_html(file_path: Path) -> str:
    document = Document(str(file_path))
    parts = ['<article class="docx-preview">']
    in_list = False

    def close_list():
        nonlocal in_list
        if in_list:
            parts.append('</ul>')
            in_list = False

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        style_name = (paragraph.style.name or '').lower()

        if not text:
            close_list()
            parts.append('<p class="empty">&nbsp;</p>')
            continue

        if 'list' in style_name:
            if not in_list:
                parts.append('<ul>')
                in_list = True
            parts.append(f'<li>{html.escape(text)}</li>')
            continue

        close_list()

        if 'heading' in style_name:
            level_digits = ''.join(filter(str.isdigit, style_name)) or '3'
            level = min(max(int(level_digits), 2), 4)
            parts.append(f'<h{level}>{html.escape(text)}</h{level}>')
        else:
            parts.append(f'<p>{html.escape(text)}</p>')

    close_list()
    parts.append('</article>')
    return ''.join(parts)


def _render_pptx_to_html(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    parts = ['<article class="pptx-preview">']
    for idx, slide in enumerate(presentation.slides, start=1):
        parts.append(f'<section class="slide"><h3>第 {idx} 页</h3>')
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            lines = []
            for paragraph in shape.text_frame.paragraphs:
                line = ''.join(run.text for run in paragraph.runs).strip()
                if line:
                    lines.append(html.escape(line))
            if lines:
                parts.append('<p>' + '<br/>'.join(lines) + '</p>')
        parts.append('</section>')
    parts.append('</article>')
    return ''.join(parts)


def _render_text_to_html(file_path: Path) -> str:
    content = file_path.read_text(encoding='utf-8', errors='ignore')
    return f'<pre class="text-preview">{html.escape(content)}</pre>'


def _load_material_entries():
    materials_path = _materials_json_path()
    if not materials_path.exists():
        raise FileNotFoundError(f'{materials_path} not found')
    with open(materials_path, 'r', encoding='utf-8') as handle:
        data = json.load(handle)
    updated_at = datetime.fromtimestamp(materials_path.stat().st_mtime, tz=dt_timezone.utc)
    return data, updated_at

class MaterialCategoryViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = MaterialCategory.objects.all().prefetch_related(
        'items__attachments'  
    )
    lookup_field = 'slug'
    filter_backends = [filters.SearchFilter, filters.OrderingFilter]
    search_fields = ['name', 'description']
    ordering_fields = ['name']
    ordering = ['name']
    
    def get_serializer_class(self):
        if self.action == 'list':
            return MaterialCategoryListSerializer
        return MaterialCategorySerializer
    
    @action(detail=False, methods=['get'])
    def overview(self, request):
        """获取课程概览统计信息"""
        categories = MaterialCategory.objects.all()
        total_items = MaterialItem.objects.count()
        total_attachments = Attachment.objects.count()
        
        return Response({
            'categories_count': categories.count(),
            'items_count': total_items,
            'attachments_count': total_attachments,
            'categories': MaterialCategoryListSerializer(
                categories, many=True, context={'request': request}
            ).data
        })

class MaterialItemViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = MaterialItem.objects.select_related('category').prefetch_related(
        'attachments'  
    )
    lookup_field = 'slug'
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    filterset_fields = ['category__slug']
    search_fields = ['title', 'description']
    ordering_fields = ['title']
    ordering = ['title']
    
    def get_serializer_class(self):
        if self.action == 'list':
            return MaterialItemListSerializer
        return MaterialItemSerializer
    
    @action(detail=True, methods=['get'])
    def attachments(self, request, slug=None):
        """获取特定项目的所有附件"""
        item = self.get_object()
        attachments = item.attachments.all()  
        serializer = AttachmentSerializer(
            attachments, many=True, context={'request': request}
        )
        return Response(serializer.data)

class AttachmentViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = Attachment.objects.select_related('material__category')
    serializer_class = AttachmentSerializer
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    filterset_fields = ['file_type', 'material__category__slug']
    search_fields = ['title', 'original_filename', 'description']
    ordering_fields = ['title', 'uploaded_at', 'download_count']
    ordering = ['-uploaded_at']
    
    @action(detail=True, methods=['get'])
    def download(self, request, pk=None):
        """下载文件"""
        attachment = self.get_object()
        
        # 更新下载计数和访问时间
        attachment.download_count += 1
        attachment.last_accessed = timezone.now()
        attachment.save(update_fields=['download_count', 'last_accessed'])
        
        # 检查文件是否存在
        if not attachment.file or not os.path.exists(attachment.file.path):
            return Response(
                {"error": "文件不存在"}, 
                status=status.HTTP_404_NOT_FOUND
            )
        
        try:
            response = FileResponse(
                open(attachment.file.path, 'rb'),
                as_attachment=True,
                filename=attachment.original_filename
            )
            
            # 设置正确的Content-Type
            content_type, _ = mimetypes.guess_type(attachment.original_filename)
            if content_type:
                response['Content-Type'] = content_type
                
            return response
            
        except Exception as e:
            return Response(
                {"error": f"文件读取失败: {str(e)}"}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    @action(detail=True, methods=['get'])
    def preview(self, request, pk=None):
        """预览文件"""
        attachment = self.get_object()
        
        if not attachment.preview_available:
            return Response(
                {"error": "此文件不支持预览"}, 
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # 更新访问时间
        attachment.last_accessed = timezone.now()
        attachment.save(update_fields=['last_accessed'])
        
        # 检查文件是否存在
        if not attachment.file or not os.path.exists(attachment.file.path):
            return Response(
                {"error": "文件不存在"}, 
                status=status.HTTP_404_NOT_FOUND
            )
        
        try:
            # 设置正确的Content-Type用于预览
            content_type, _ = mimetypes.guess_type(attachment.original_filename)
            
            response = FileResponse(
                open(attachment.file.path, 'rb'),
                as_attachment=False  # 不作为附件，让浏览器尝试预览
            )
            
            if content_type:
                response['Content-Type'] = content_type
            
            # 对于PDF文件或其他文件类型，设置inline显示
            if content_type == 'application/pdf':
                response['Content-Disposition'] = f'inline; filename="{attachment.original_filename}"'
            
            return response
            
        except Exception as e:
            return Response(
                {"error": f"文件读取失败: {str(e)}"}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
    @action(detail=False, methods=['get'])
    def statistics(self, request):
        """获取文件统计信息"""
        attachments = self.get_queryset()
        
        # 按文件类型统计
        type_stats = {}
        for attachment in attachments:
            file_type = attachment.get_file_type_display()
            if file_type not in type_stats:
                type_stats[file_type] = {'count': 0, 'total_size': 0}
            type_stats[file_type]['count'] += 1
            type_stats[file_type]['total_size'] += attachment.file_size or 0
        
        return Response({
            'total_count': attachments.count(),
            'total_size': sum(a.file_size or 0 for a in attachments),
            'type_statistics': type_stats,
            'most_downloaded': AttachmentSerializer(
                attachments.order_by('-download_count')[:5],
                many=True,
                context={'request': request}
            ).data
        })

    @action(detail=False, methods=['get'], url_path='by-experiment')
    def by_experiment(self, request):
        category = request.query_params.get('category')
        if not category:
            return Response({"detail": "category query parameter is required."},
                            status=status.HTTP_400_BAD_REQUEST)

        attachments = (
            self.get_queryset()
            .filter(material__category__name=category)
            .order_by('original_filename')
        )
        serializer = self.get_serializer(attachments, many=True)
        material = MaterialItem.objects.filter(category__name=category).first()

        return Response({
            "category": category,
            "material": MaterialItemSerializer(material).data if material else None,
            "attachments": serializer.data,
        })

    @action(detail=False, methods=['get'], url_path='books')
    def books(self, request):
        category_names = request.query_params.getlist('category') or ['参考书籍', '参考资料']
        attachments = (
            self.get_queryset()
            .filter(material__category__name__in=category_names)
            .order_by('material__category__name', 'original_filename')
        )
        serializer = self.get_serializer(attachments, many=True)
        return Response({"categories": category_names, "attachments": serializer.data})


class MaterialCatalogView(APIView):
    permission_classes = [permissions.AllowAny]

    def get(self, request):
        try:
            entries, updated_at = _load_material_entries()
        except FileNotFoundError:
            return Response(
                {"detail": "materials.json not found"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )
        except json.JSONDecodeError as exc:
            return Response(
                {"detail": f"materials.json parse error: {exc}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

        experiments: dict[str, dict] = {}
        videos = []
        books = []
        latest_order = -1
        latest_label: Optional[str] = None

        for entry in entries:
            try:
                attachment = _build_attachment_payload(entry, request)
            except (ValueError, FileNotFoundError):
                # Skip malformed entries quietly so the rest of the data remains accessible.
                continue

            category = attachment['category']
            order = attachment['category_order']

            if category in BOOK_CATEGORIES:
                books.append(attachment)
                continue

            if attachment['file_type'] in VIDEO_FILE_TYPES:
                videos.append(attachment)

            if order >= 1000:
                continue

            bucket = experiments.setdefault(
                category,
                {
                    'category': category,
                    'category_display': attachment['category_display'],
                    'order': order,
                    'items': [],
                },
            )
            bucket['items'].append(attachment)
            bucket['files_count'] = len(bucket['items'])

            if order > latest_order:
                latest_order = order
                latest_label = attachment['category_display']

        experiment_list = sorted(
            experiments.values(), key=lambda item: (item['order'], item['category_display'])
        )
        for bucket in experiment_list:
            bucket['items'].sort(
                key=lambda att: ((att.get('item_label') or ''), att['label'])
            )

        videos.sort(key=lambda att: (att.get('category_order', 999), att['label']))

        def _book_category_index(name: str) -> int:
            return BOOK_CATEGORY_ORDER.index(name) if name in BOOK_CATEGORY_ORDER else len(BOOK_CATEGORY_ORDER)

        books.sort(key=lambda att: (_book_category_index(att['category']), att['label']))

        return Response(
            {
                'updated_at': updated_at.isoformat() if updated_at else None,
                'synced_to': {
                    'order': latest_order if latest_order >= 0 else None,
                    'label': latest_label,
                },
                'experiments': experiment_list,
                'videos': videos,
                'books': books,
            }
        )


class BaseMaterialFileView(APIView):
    permission_classes = [permissions.AllowAny]
    as_attachment = False

    def get(self, request):
        relative_path = request.query_params.get('path')
        if not relative_path:
            return Response(
                {"detail": "path query parameter is required."},
                status=status.HTTP_400_BAD_REQUEST,
            )
        try:
            storage_relative, file_path = _resolve_material_file(relative_path)
        except (ValueError, FileNotFoundError):
            return Response(
                {"detail": "文件不存在或路径非法"}, status=status.HTTP_404_NOT_FOUND
            )

        try:
            file_handle = open(file_path, 'rb')
        except OSError as exc:
            raise Http404(str(exc)) from exc

        response = FileResponse(
            file_handle,
            as_attachment=self.as_attachment,
            filename=file_path.name,
        )
        content_type, _ = mimetypes.guess_type(file_path.name)
        if content_type:
            response['Content-Type'] = content_type

        disposition = 'attachment' if self.as_attachment else 'inline'
        encoded_name = quote(file_path.name)
        response['Content-Disposition'] = f'{disposition}; filename="{encoded_name}"'
        if not self.as_attachment:
            response['X-Frame-Options'] = 'ALLOWALL'
        return response


class MaterialFilePreviewView(BaseMaterialFileView):
    as_attachment = False


class MaterialFileDownloadView(BaseMaterialFileView):
    as_attachment = True


class MaterialFileHtmlPreviewView(APIView):
    permission_classes = [permissions.AllowAny]

    def get(self, request):
        relative_path = request.query_params.get('path')
        if not relative_path:
            return Response({"detail": "path query parameter is required."}, status=status.HTTP_400_BAD_REQUEST)

        try:
            _, file_path = _resolve_material_file(relative_path)
        except (ValueError, FileNotFoundError):
            return Response({"detail": "文件不存在或路径非法"}, status=status.HTTP_404_NOT_FOUND)

        suffix = file_path.suffix.lower()
        try:
            if suffix in DOC_EXTENSIONS:
                html_payload = _render_docx_to_html(file_path)
            elif suffix in PPT_EXTENSIONS:
                html_payload = _render_pptx_to_html(file_path)
            elif suffix in TEXT_EXTENSIONS:
                html_payload = _render_text_to_html(file_path)
            else:
                return Response({"detail": "此文件类型暂不支持 HTML 预览"}, status=status.HTTP_400_BAD_REQUEST)
        except Exception as exc:  # pragma: no cover - conversion edge cases
            return Response({"detail": f"预览生成失败: {exc}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        return Response({"html": html_payload})